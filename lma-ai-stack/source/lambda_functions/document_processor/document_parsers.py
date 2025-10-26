#!/usr/bin/env python3
# Copyright (c) 2025
# This file is licensed under the MIT License.

"""
Enhanced Document Parsers for RAG Knowledge Base
Integrated RAGFlow features: Context-aware chunking, Table tracking, Embedded files
Supports: PDF, DOCX, PPTX, XLSX, TXT, MD, HTML, JSON, CSV, Code files

ALL FEATURES ARE LIGHTWEIGHT - NO GPU/ML MODELS REQUIRED
"""

import os
import re
import io
import json
import csv
import logging
import hashlib
import zipfile
from typing import List, Dict, Any, Optional, Tuple
from pathlib import Path
from io import BytesIO

logger = logging.getLogger(__name__)
logger.setLevel(logging.INFO)


def num_tokens_from_string(text: str) -> int:
    """Estimate token count (rough approximation: 1 token ≈ 4 characters)"""
    return len(text) // 4


# ==================== EMBEDDED FILE EXTRACTION ====================
# Integrated from RAGFlow - Extract files within files (DOCX, XLSX, OLE)

def _is_zip(h: bytes) -> bool:
    """Check if bytes are ZIP format"""
    return h.startswith(b"PK\x03\x04") or h.startswith(b"PK\x05\x06") or h.startswith(b"PK\x07\x08")


def _is_pdf(h: bytes) -> bool:
    """Check if bytes are PDF format"""
    return h.startswith(b"%PDF-")


def _is_ole(h: bytes) -> bool:
    """Check if bytes are OLE format (old MS Office)"""
    return h.startswith(b"\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1")


def _sha10(b: bytes) -> str:
    """Generate short hash for file identification"""
    return hashlib.sha256(b).hexdigest()[:10]


def _guess_ext(b: bytes) -> str:
    """Guess file extension from binary content"""
    h = b[:8]
    if _is_zip(h):
        try:
            with zipfile.ZipFile(BytesIO(b), "r") as z:
                names = [n.lower() for n in z.namelist()]
                if any(n.startswith("word/") for n in names):
                    return ".docx"
                if any(n.startswith("ppt/") for n in names):
                    return ".pptx"
                if any(n.startswith("xl/") for n in names):
                    return ".xlsx"
        except Exception:
            pass
        return ".zip"
    if _is_pdf(h):
        return ".pdf"
    if _is_ole(h):
        return ".doc"
    return ".bin"


def extract_embedded_files(file_content: bytes, max_file_size: int = 50 * 1024 * 1024) -> List[Tuple[str, bytes]]:
    """
    Extract embedded files from DOCX/XLSX/PPTX (first layer only)
    
    This is a LIGHTWEIGHT operation - just ZIP extraction, no ML
    
    Args:
        file_content: Binary content of file
        max_file_size: Maximum size of embedded file to extract (default 50MB)
    
    Returns:
        List of (filename, file_bytes) tuples
    """
    # Safety check: Don't process files too large
    if len(file_content) > max_file_size:
        logger.warning(f"File too large for embedded extraction: {len(file_content)} bytes")
        return []
    
    head = file_content[:8]
    result: List[Tuple[str, bytes]] = []
    seen = set()
    
    def add_file(b: bytes, name_hint: str = ""):
        """Add unique file to result list"""
        # Skip files that are too large (avoid memory issues)
        if len(b) > max_file_size:
            logger.warning(f"Skipping large embedded file: {len(b)} bytes")
            return
        
        # Skip empty files
        if len(b) == 0:
            return
        
        h10 = _sha10(b)
        if h10 in seen:
            return
        seen.add(h10)
        
        ext = _guess_ext(b)
        if "." in name_hint:
            fname = name_hint.split("/")[-1]
        else:
            fname = f"embedded_{h10}{ext}"
        result.append((fname, b))
    
    # OOXML/ZIP container (docx/xlsx/pptx)
    if _is_zip(head):
        try:
            with zipfile.ZipFile(BytesIO(file_content), "r") as z:
                embed_dirs = (
                    "word/embeddings/", "word/objects/", "word/activex/",
                    "xl/embeddings/", "ppt/embeddings/"
                )
                for name in z.namelist():
                    low = name.lower()
                    if any(low.startswith(d) for d in embed_dirs):
                        try:
                            b = z.read(name)
                            add_file(b, name)
                        except Exception:
                            pass
        except Exception:
            pass
    
    # OLE container (old doc/xls/ppt) - requires olefile
    if _is_ole(head):
        try:
            import olefile
            with olefile.OleFileIO(BytesIO(file_content)) as ole:
                for entry in ole.listdir():
                    p = "/".join(entry)
                    try:
                        data = ole.openstream(entry).read()
                        if data and ("Ole10Native" in p or "ole10native" in p.lower()):
                            # Extract embedded payload
                            add_file(data, p)
                    except Exception:
                        continue
        except ImportError:
            logger.warning("olefile not available, skipping OLE embedded file extraction")
        except Exception:
            pass
    
    return result


class TextParser:
    """Parser for plain text files"""
    
    @staticmethod
    def parse(file_content: bytes, filename: str) -> str:
        """Parse text file"""
        try:
            # Try UTF-8 first
            return file_content.decode('utf-8')
        except UnicodeDecodeError:
            # Fallback to latin-1
            return file_content.decode('latin-1', errors='ignore')


class MarkdownParser:
    """Parser for Markdown files"""
    
    @staticmethod
    def parse(file_content: bytes, filename: str) -> str:
        """Parse markdown file"""
        try:
            text = file_content.decode('utf-8')
            # Remove code blocks for cleaner indexing
            text = re.sub(r'```[\s\S]*?```', '[CODE BLOCK]', text)
            # Remove inline code
            text = re.sub(r'`[^`]+`', '[CODE]', text)
            # Remove images
            text = re.sub(r'!\[([^\]]*)\]\([^\)]+\)', r'[Image: \1]', text)
            return text
        except Exception as e:
            logger.error(f"Error parsing markdown: {str(e)}")
            return file_content.decode('utf-8', errors='ignore')


class PDFParser:
    """Parser for PDF files using PyPDF2"""
    
    @staticmethod
    def parse(file_content: bytes, filename: str) -> str:
        """Parse PDF file"""
        try:
            import PyPDF2
            
            pdf_file = io.BytesIO(file_content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            text_parts = []
            for page_num, page in enumerate(pdf_reader.pages):
                text = page.extract_text()
                if text.strip():
                    text_parts.append(f"[Page {page_num + 1}]\n{text}")
            
            return "\n\n".join(text_parts)
        
        except Exception as e:
            logger.error(f"Error parsing PDF with PyPDF2: {str(e)}")
            # Fallback to pdfplumber if available
            try:
                import pdfplumber
                
                pdf_file = io.BytesIO(file_content)
                text_parts = []
                
                with pdfplumber.open(pdf_file) as pdf:
                    for page_num, page in enumerate(pdf.pages):
                        text = page.extract_text()
                        if text:
                            text_parts.append(f"[Page {page_num + 1}]\n{text}")
                
                return "\n\n".join(text_parts)
            
            except Exception as e2:
                logger.error(f"Error parsing PDF with pdfplumber: {str(e2)}")
                return f"Error parsing PDF: {str(e2)}"


class DOCXParser:
    """
    Enhanced DOCX Parser with Table Location Tracking
    Integrated from RAGFlow - tracks heading hierarchy for tables
    """
    
    @staticmethod
    def parse(file_content: bytes, filename: str) -> str:
        """
        Parse DOCX file with context preservation
        
        LIGHTWEIGHT: No ML, just text parsing and structure tracking
        """
        try:
            from docx import Document
            from docx.text.paragraph import Paragraph
            
            docx_file = io.BytesIO(file_content)
            doc = Document(docx_file)
            
            text_parts = []
            
            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
            
            # Extract tables with location context
            for table_idx, table in enumerate(doc.tables):
                # Get nearest heading for this table (RAGFlow feature)
                table_location = DOCXParser._get_table_location(
                    doc, table_idx, filename
                )
                
                table_text = f"\n[Table {table_idx + 1}"
                if table_location:
                    table_text += f" - Location: {table_location}"
                table_text += "]\n"
                
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells])
                    if row_text.strip():
                        table_text += row_text + "\n"
                text_parts.append(table_text)
            
            return "\n\n".join(text_parts)
        
        except Exception as e:
            logger.error(f"Error parsing DOCX: {str(e)}")
            return f"Error parsing DOCX: {str(e)}"
    
    @staticmethod
    def _get_table_location(doc, table_index: int, filename: str) -> str:
        """
        Get hierarchical location of table (Document > Section > Subsection)
        
        This is LIGHTWEIGHT - just text parsing, no ML
        Integrated from RAGFlow's __get_nearest_title method
        
        Args:
            doc: Document object
            table_index: Index of table
            filename: Document filename
        
        Returns:
            Location string like "DocName > Section 1 > Table"
        """
        try:
            from docx.text.paragraph import Paragraph
            
            # Get document name from filename
            doc_name = re.sub(r"\.[a-zA-Z]+$", "", filename)
            if not doc_name:
                doc_name = "Document"
            
            # Build list of all blocks (paragraphs and tables)
            blocks = []
            for i, block in enumerate(doc._element.body):
                if block.tag.endswith('p'):  # Paragraph
                    p = Paragraph(block, doc)
                    blocks.append(('p', i, p))
                elif block.tag.endswith('tbl'):  # Table
                    blocks.append(('t', i, None))
            
            # Find target table position
            target_pos = -1
            table_count = 0
            for i, (block_type, pos, _) in enumerate(blocks):
                if block_type == 't':
                    if table_count == table_index:
                        target_pos = pos
                        break
                    table_count += 1
            
            if target_pos == -1:
                return ""
            
            # Find nearest heading before table
            titles = []
            current_level = 999
            max_search_blocks = 100  # Limit search to prevent performance issues
            blocks_searched = 0
            
            for i in range(len(blocks) - 1, -1, -1):
                block_type, pos, block = blocks[i]
                
                if pos >= target_pos:  # Skip blocks after table
                    continue
                
                if block_type != 'p':
                    continue
                
                blocks_searched += 1
                if blocks_searched > max_search_blocks:
                    logger.warning(f"Table location search limit reached ({max_search_blocks} blocks)")
                    break
                
                # Check if paragraph is a heading
                if block.style and block.style.name:
                    heading_match = re.search(r"Heading\s*(\d+)", block.style.name, re.I)
                    if heading_match:
                        level = int(heading_match.group(1))
                        title_text = block.text.strip()
                        
                        if title_text and level <= 7:
                            if level < current_level:
                                titles.insert(0, title_text)
                                current_level = level
                                
                                # Stop if we found a top-level heading
                                if level == 1:
                                    break
            
            # Build hierarchy string
            if titles:
                hierarchy = [doc_name] + titles
                return " > ".join(hierarchy)
            
            return doc_name
        
        except Exception as e:
            logger.warning(f"Error getting table location: {e}")
            return ""


class PPTXParser:
    """Parser for PPTX files using python-pptx"""
    
    @staticmethod
    def parse(file_content: bytes, filename: str) -> str:
        """Parse PPTX file"""
        try:
            from pptx import Presentation
            
            pptx_file = io.BytesIO(file_content)
            prs = Presentation(pptx_file)
            
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = f"[Slide {slide_num + 1}]\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                    
                    # Extract table content
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = " | ".join([cell.text.strip() for cell in row.cells])
                            if row_text.strip():
                                slide_text += row_text + "\n"
                
                if slide_text.strip() != f"[Slide {slide_num + 1}]":
                    text_parts.append(slide_text)
            
            return "\n\n".join(text_parts)
        
        except Exception as e:
            logger.error(f"Error parsing PPTX: {str(e)}")
            return f"Error parsing PPTX: {str(e)}"


class XLSXParser:
    """Parser for XLSX files using openpyxl"""
    
    @staticmethod
    def parse(file_content: bytes, filename: str) -> str:
        """Parse XLSX file"""
        try:
            from openpyxl import load_workbook
            
            xlsx_file = io.BytesIO(file_content)
            workbook = load_workbook(xlsx_file, data_only=True)
            
            text_parts = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_text = f"[Sheet: {sheet_name}]\n"
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        sheet_text += row_text + "\n"
                
                if sheet_text.strip() != f"[Sheet: {sheet_name}]":
                    text_parts.append(sheet_text)
            
            return "\n\n".join(text_parts)
        
        except Exception as e:
            logger.error(f"Error parsing XLSX: {str(e)}")
            return f"Error parsing XLSX: {str(e)}"


class HTMLParser:
    """Parser for HTML files - LIGHTWEIGHT"""
    
    @staticmethod
    def parse(file_content: bytes, filename: str) -> str:
        """Parse HTML file"""
        try:
            from bs4 import BeautifulSoup
            
            html = file_content.decode('utf-8', errors='ignore')
            soup = BeautifulSoup(html, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            # Get text
            text = soup.get_text()
            
            # Clean up whitespace
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = '\n'.join(chunk for chunk in chunks if chunk)
            
            return text
        
        except ImportError:
            # Fallback: Simple regex-based HTML stripping
            html = file_content.decode('utf-8', errors='ignore')
            text = re.sub(r'<script[^>]*>.*?</script>', '', html, flags=re.DOTALL)
            text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL)
            text = re.sub(r'<[^>]+>', ' ', text)
            text = re.sub(r'\s+', ' ', text)
            return text.strip()
        
        except Exception as e:
            logger.error(f"Error parsing HTML: {str(e)}")
            return f"Error parsing HTML: {str(e)}"


class JSONParser:
    """Parser for JSON/JSONL files - LIGHTWEIGHT"""
    
    @staticmethod
    def parse(file_content: bytes, filename: str) -> str:
        """Parse JSON file"""
        try:
            text = file_content.decode('utf-8', errors='ignore')
            
            # Handle JSONL (JSON Lines) format
            if filename.endswith(('.jsonl', '.ldjson')):
                lines = text.strip().split('\n')
                parsed = []
                for line in lines:
                    if line.strip():
                        try:
                            obj = json.loads(line)
                            parsed.append(json.dumps(obj, indent=2, ensure_ascii=False))
                        except:
                            parsed.append(line)
                return "\n\n".join(parsed)
            
            # Regular JSON
            data = json.loads(text)
            return json.dumps(data, indent=2, ensure_ascii=False)
        
        except Exception as e:
            logger.error(f"Error parsing JSON: {str(e)}")
            # Return raw text as fallback
            return file_content.decode('utf-8', errors='ignore')


class CSVParser:
    """Parser for CSV files - LIGHTWEIGHT"""
    
    @staticmethod
    def parse(file_content: bytes, filename: str) -> str:
        """Parse CSV file"""
        try:
            text = file_content.decode('utf-8', errors='ignore')
            
            # Parse CSV
            reader = csv.reader(io.StringIO(text))
            rows = list(reader)
            
            if not rows:
                return ""
            
            # Convert to markdown table
            result = []
            for i, row in enumerate(rows):
                result.append(" | ".join(row))
                # Add separator after header
                if i == 0:
                    result.append(" | ".join(["---"] * len(row)))
            
            return "\n".join(result)
        
        except Exception as e:
            logger.error(f"Error parsing CSV: {str(e)}")
            return file_content.decode('utf-8', errors='ignore')


class CodeParser:
    """Parser for code files (py, js, java, etc.) - LIGHTWEIGHT"""
    
    @staticmethod
    def parse(file_content: bytes, filename: str) -> str:
        """Parse code file"""
        try:
            # Try UTF-8 first
            text = file_content.decode('utf-8')
        except UnicodeDecodeError:
            # Fallback to latin-1
            text = file_content.decode('latin-1', errors='ignore')
        
        # Add file extension info for better context
        ext = Path(filename).suffix.lower()
        lang_map = {
            '.py': 'Python',
            '.js': 'JavaScript',
            '.ts': 'TypeScript',
            '.java': 'Java',
            '.c': 'C',
            '.cpp': 'C++',
            '.h': 'C/C++ Header',
            '.go': 'Go',
            '.rs': 'Rust',
            '.rb': 'Ruby',
            '.php': 'PHP',
            '.sh': 'Shell',
            '.sql': 'SQL',
            '.kt': 'Kotlin',
            '.cs': 'C#'
        }
        
        lang = lang_map.get(ext, 'Code')
        return f"[{lang} Code: {filename}]\n\n{text}"


class DocumentParserFactory:
    """Factory to get the appropriate parser for a file type"""
    
    PARSERS = {
        # Documents
        '.pdf': PDFParser,
        '.docx': DOCXParser,
        '.pptx': PPTXParser,
        '.xlsx': XLSXParser,
        
        # Text formats
        '.txt': TextParser,
        '.md': MarkdownParser,
        '.markdown': MarkdownParser,
        
        # Structured data
        '.html': HTMLParser,
        '.htm': HTMLParser,
        '.json': JSONParser,
        '.jsonl': JSONParser,
        '.ldjson': JSONParser,
        '.csv': CSVParser,
        
        # Code files
        '.py': CodeParser,
        '.js': CodeParser,
        '.ts': CodeParser,
        '.java': CodeParser,
        '.c': CodeParser,
        '.cpp': CodeParser,
        '.h': CodeParser,
        '.go': CodeParser,
        '.rs': CodeParser,
        '.rb': CodeParser,
        '.php': CodeParser,
        '.sh': CodeParser,
        '.sql': CodeParser,
        '.kt': CodeParser,
        '.cs': CodeParser,
    }
    
    @classmethod
    def get_parser(cls, filename: str):
        """Get parser based on file extension"""
        ext = Path(filename).suffix.lower()
        parser = cls.PARSERS.get(ext)
        
        if not parser:
            raise ValueError(f"Unsupported file type: {ext}")
        
        return parser
    
    @classmethod
    def parse_document(cls, file_content: bytes, filename: str) -> str:
        """Parse document and return extracted text"""
        parser = cls.get_parser(filename)
        return parser.parse(file_content, filename)


class TextChunker:
    """
    Enhanced Text Chunker with Context Preservation
    Integrated from RAGFlow - preserves heading hierarchy in chunks
    
    LIGHTWEIGHT: No ML, just smart text processing
    """
    
    def __init__(self, chunk_size: int = 512, overlap: int = 50):
        """
        Initialize chunker
        
        Args:
            chunk_size: Target chunk size in tokens
            overlap: Number of tokens to overlap between chunks
        """
        self.chunk_size = chunk_size
        self.overlap = overlap
    
    def extract_heading_hierarchy(self, text: str) -> Dict[int, str]:
        """
        Extract heading hierarchy from text
        
        LIGHTWEIGHT: Just regex parsing
        
        Args:
            text: Full document text
        
        Returns:
            Dict mapping line number to heading hierarchy
        """
        hierarchy = {}
        current_headings = {}  # level -> heading text
        
        lines = text.split('\n')
        
        for i, line in enumerate(lines):
            # Check for markdown headings
            md_heading = re.match(r'^(#{1,6})\s+(.+)', line)
            if md_heading:
                level = len(md_heading.group(1))
                heading = md_heading.group(2).strip()
                
                # Update current headings at this level
                current_headings[level] = heading
                
                # Clear deeper levels
                keys_to_remove = [k for k in current_headings.keys() if k > level]
                for k in keys_to_remove:
                    del current_headings[k]
                
                # Build hierarchy path
                path_parts = [current_headings[l] for l in sorted(current_headings.keys())]
                hierarchy[i] = " > ".join(path_parts)
            
            # Check for [Page X] or [Table X] markers
            elif line.strip().startswith('[') and ']' in line:
                marker = line.strip()
                if marker not in current_headings.values():
                    hierarchy[i] = marker
        
        return hierarchy
    
    def chunk_by_sentences(self, text: str) -> List[str]:
        """
        Chunk text by sentences with overlap
        
        Args:
            text: Text to chunk
        
        Returns:
            List of text chunks
        """
        # Split into sentences
        sentence_pattern = r'(?<=[.!?])\s+'
        sentences = re.split(sentence_pattern, text)
        
        chunks = []
        current_chunk = []
        current_tokens = 0
        
        for sentence in sentences:
            sentence_tokens = num_tokens_from_string(sentence)
            
            # If adding this sentence exceeds chunk size, save current chunk
            if current_tokens + sentence_tokens > self.chunk_size and current_chunk:
                chunks.append(" ".join(current_chunk))
                
                # Keep last few sentences for overlap
                overlap_sentences = []
                overlap_tokens = 0
                for s in reversed(current_chunk):
                    s_tokens = num_tokens_from_string(s)
                    if overlap_tokens + s_tokens <= self.overlap:
                        overlap_sentences.insert(0, s)
                        overlap_tokens += s_tokens
                    else:
                        break
                
                current_chunk = overlap_sentences
                current_tokens = overlap_tokens
            
            current_chunk.append(sentence)
            current_tokens += sentence_tokens
        
        # Add final chunk
        if current_chunk:
            chunks.append(" ".join(current_chunk))
        
        return chunks
    
    def chunk_by_paragraphs(self, text: str) -> List[str]:
        """
        Chunk text by paragraphs with size limits
        
        Args:
            text: Text to chunk
        
        Returns:
            List of text chunks
        """
        paragraphs = text.split('\n\n')
        chunks = []
        current_chunk = []
        current_tokens = 0
        
        for para in paragraphs:
            para = para.strip()
            if not para:
                continue
            
            para_tokens = num_tokens_from_string(para)
            
            # If paragraph itself is too large, split it by sentences
            if para_tokens > self.chunk_size:
                # Save current chunk first
                if current_chunk:
                    chunks.append("\n\n".join(current_chunk))
                    current_chunk = []
                    current_tokens = 0
                
                # Split large paragraph
                para_chunks = self.chunk_by_sentences(para)
                chunks.extend(para_chunks)
                continue
            
            # If adding this paragraph exceeds chunk size, save current chunk
            if current_tokens + para_tokens > self.chunk_size and current_chunk:
                chunks.append("\n\n".join(current_chunk))
                current_chunk = [para]
                current_tokens = para_tokens
            else:
                current_chunk.append(para)
                current_tokens += para_tokens
        
        # Add final chunk
        if current_chunk:
            chunks.append("\n\n".join(current_chunk))
        
        return chunks
    
    def chunk_text(self, text: str, strategy: str = 'paragraphs') -> List[str]:
        """
        Chunk text using specified strategy
        
        Args:
            text: Text to chunk
            strategy: 'sentences' or 'paragraphs'
        
        Returns:
            List of text chunks
        """
        if strategy == 'sentences':
            return self.chunk_by_sentences(text)
        else:
            return self.chunk_by_paragraphs(text)
    
    def chunk_text_with_context(
        self, 
        text: str, 
        filename: str = "Document",
        strategy: str = 'paragraphs'
    ) -> List[Dict[str, Any]]:
        """
        Chunk text with context preservation (RAGFlow feature)
        
        LIGHTWEIGHT: Just text processing
        
        Args:
            text: Text to chunk
            filename: Document filename for context
            strategy: 'sentences' or 'paragraphs'
        
        Returns:
            List of dicts with 'content', 'context', and 'enhanced_text'
        """
        # Safety check: Don't process empty or whitespace-only text
        if not text or not text.strip():
            logger.warning("Empty or whitespace-only text provided for chunking")
            return []
        
        # Extract heading hierarchy
        hierarchy = self.extract_heading_hierarchy(text)
        
        # Get basic chunks
        chunks = self.chunk_text(text, strategy)
        
        # Filter out empty chunks
        chunks = [c for c in chunks if c and c.strip()]
        
        # Enhance chunks with context
        result = []
        doc_name = re.sub(r"\.[a-zA-Z]+$", "", filename)
        
        for i, chunk in enumerate(chunks):
            # Find the nearest heading for this chunk
            chunk_lines = chunk.split('\n')
            context = doc_name
            
            # Try to find heading context from chunk content
            for line in chunk_lines:
                # Check if line starts a heading
                if line.strip().startswith('#'):
                    heading = line.strip('#').strip()
                    context = f"{doc_name} > {heading}"
                    break
                # Check for markers
                elif line.strip().startswith('[') and ']' in line:
                    marker = line.strip()
                    context = f"{doc_name} > {marker}"
                    break
            
            # Build enhanced text with context
            enhanced_text = f"""[Context: {context}]

{chunk}"""
            
            result.append({
                'content': chunk,
                'context': context,
                'enhanced_text': enhanced_text,
                'chunk_index': i,
                'total_chunks': len(chunks)
            })
        
        return result


# For testing
if __name__ == "__main__":
    # Test text parser
    test_text = "This is a test document.\n\nIt has multiple paragraphs.\n\nEach paragraph is separated by blank lines."
    
    chunker = TextChunker(chunk_size=20, overlap=5)
    chunks = chunker.chunk_text(test_text, strategy='paragraphs')
    
    print(f"Original text length: {len(test_text)}")
    print(f"Number of chunks: {len(chunks)}")
    for i, chunk in enumerate(chunks):
        print(f"\nChunk {i+1}:")
        print(chunk)


